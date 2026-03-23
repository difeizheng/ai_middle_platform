"""
文档解析服务
支持 PDF、Word、Excel、PPT、TXT、Markdown 等格式
"""
import os
from typing import Optional, Dict, Any, List
from pathlib import Path
import hashlib

from ..core.logger import get_logger

logger = get_logger(__name__)


class DocumentParser:
    """
    文档解析器
    """

    def __init__(self):
        self.supported_formats = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".xlsx": self._parse_xlsx,
            ".pptx": self._parse_pptx,
            ".txt": self._parse_txt,
            ".md": self._parse_md,
            ".csv": self._parse_csv,
        }

    def parse(self, file_path: str) -> Dict[str, Any]:
        """
        解析文档

        Args:
            file_path: 文件路径

        Returns:
            {
                "content": str,  # 完整文本内容
                "metadata": dict,  # 元数据
                "pages": list,  # 分页内容（如果有）
            }
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在：{file_path}")

        # 获取文件扩展名
        ext = Path(file_path).suffix.lower()

        if ext not in self.supported_formats:
            raise ValueError(f"不支持的文件格式：{ext}")

        # 计算文件 hash
        file_hash = self._calculate_hash(file_path)

        # 获取文件大小
        file_size = os.path.getsize(file_path)

        # 解析文件
        parse_func = self.supported_formats[ext]
        result = parse_func(file_path)

        # 添加元数据
        result["metadata"]["file_path"] = file_path
        result["metadata"]["file_name"] = os.path.basename(file_path)
        result["metadata"]["file_hash"] = file_hash
        result["metadata"]["file_size"] = file_size
        result["metadata"]["file_type"] = ext[1:]

        logger.info(f"文档解析成功：{file_path}, 内容长度：{len(result['content'])}")

        return result

    def _calculate_hash(self, file_path: str) -> str:
        """计算文件 SHA256 哈希"""
        sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                sha256.update(chunk)
        return sha256.hexdigest()

    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """解析 PDF 文件"""
        try:
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            pages = []
            content_parts = []

            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                pages.append({
                    "page_num": i + 1,
                    "content": text,
                })
                content_parts.append(text)

            return {
                "content": "\n\n".join(content_parts),
                "metadata": {
                    "page_count": len(reader.pages),
                    "format": "pdf",
                },
                "pages": pages,
            }
        except ImportError:
            logger.warning("pypdf 未安装，PDF 解析功能不可用")
            return {
                "content": "",
                "metadata": {"format": "pdf", "error": "缺少依赖：pypdf"},
                "pages": [],
            }
        except Exception as e:
            logger.error(f"PDF 解析失败：{e}")
            raise

    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """解析 Word 文件"""
        try:
            from docx import Document

            doc = Document(file_path)
            paragraphs = []
            tables = []

            # 提取段落
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)

            # 提取表格
            for i, table in enumerate(doc.tables):
                table_text = []
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                tables.append({
                    "table_index": i,
                    "content": "\n".join(table_text),
                })

            content = "\n\n".join(paragraphs)
            if tables:
                table_contents = ["\n".join(t["content"] for t in tables)]
                content += "\n\n[表格内容]\n\n" + "\n\n".join(table_contents)

            return {
                "content": content,
                "metadata": {
                    "paragraph_count": len(paragraphs),
                    "table_count": len(tables),
                    "format": "docx",
                },
                "pages": [],
            }
        except ImportError:
            logger.warning("python-docx 未安装")
            return {
                "content": "",
                "metadata": {"format": "docx", "error": "缺少依赖：python-docx"},
                "pages": [],
            }
        except Exception as e:
            logger.error(f"Word 解析失败：{e}")
            raise

    def _parse_xlsx(self, file_path: str) -> Dict[str, Any]:
        """解析 Excel 文件"""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, read_only=True)
            sheets = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_data):  # 跳过空行
                        rows.append(", ".join(row_data))

                sheets.append({
                    "sheet_name": sheet_name,
                    "content": "\n".join(rows),
                    "row_count": len(rows),
                })

            content = "\n\n".join([f"[{s['sheet_name']}]\n{s['content']}" for s in sheets])

            return {
                "content": content,
                "metadata": {
                    "sheet_count": len(sheets),
                    "format": "xlsx",
                },
                "pages": sheets,
            }
        except ImportError:
            logger.warning("openpyxl 未安装")
            return {
                "content": "",
                "metadata": {"format": "xlsx", "error": "缺少依赖：openpyxl"},
                "pages": [],
            }
        except Exception as e:
            logger.error(f"Excel 解析失败：{e}")
            raise

    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """解析 PPT 文件"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            slides = []

            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                slides.append({
                    "slide_num": i + 1,
                    "content": "\n".join(slide_text),
                })

            content = "\n\n".join([s["content"] for s in slides])

            return {
                "content": content,
                "metadata": {
                    "slide_count": len(slides),
                    "format": "pptx",
                },
                "pages": slides,
            }
        except ImportError:
            logger.warning("python-pptx 未安装")
            return {
                "content": "",
                "metadata": {"format": "pptx", "error": "缺少依赖：python-pptx"},
                "pages": [],
            }
        except Exception as e:
            logger.error(f"PPT 解析失败：{e}")
            raise

    def _parse_txt(self, file_path: str) -> Dict[str, Any]:
        """解析 TXT 文件"""
        try:
            # 尝试不同的编码
            encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
            content = ""

            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        content = f.read()
                    logger.info(f"使用 {encoding} 编码成功读取文件")
                    break
                except UnicodeDecodeError:
                    continue

            if not content:
                raise ValueError("无法识别文件编码")

            lines = content.split("\n")

            return {
                "content": content,
                "metadata": {
                    "line_count": len(lines),
                    "format": "txt",
                },
                "pages": [],
            }
        except Exception as e:
            logger.error(f"TXT 解析失败：{e}")
            raise

    def _parse_md(self, file_path: str) -> Dict[str, Any]:
        """解析 Markdown 文件"""
        # Markdown 作为纯文本处理
        return self._parse_txt(file_path)

    def _parse_csv(self, file_path: str) -> Dict[str, Any]:
        """解析 CSV 文件"""
        import csv

        rows = []
        with open(file_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(", ".join(row))

        content = "\n".join(rows)

        return {
            "content": content,
            "metadata": {
                "row_count": len(rows),
                "format": "csv",
            },
            "pages": [],
        }
