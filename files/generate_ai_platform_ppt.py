#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
AI 中台 PPT 生成工具
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


class AIPPTGenerator:
    """AI 中台 PPT 生成器"""

    def __init__(self, title="AI 中台系统", subtitle="企业级 AI 能力基础设施"):
        self.title = title
        self.subtitle = subtitle
        self.prs = Presentation()

        # 设置宽屏 16:9
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, title=None, subtitle=None):
        """添加标题页"""
        layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)

        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title or self.title
        subtitle_shape.text = subtitle or self.subtitle

        return slide

    def add_content_slide(self, title, content, bullets=None):
        """添加内容页"""
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)

        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = title

        tf = body_shape.text_frame
        tf.word_wrap = True

        if content:
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(18)

        if bullets:
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(16)

        return slide

    def add_section_slide(self, section_title):
        """添加章节页"""
        layout = self.prs.slide_layouts[2]
        slide = self.prs.slides.add_slide(layout)

        title_shape = slide.shapes.title
        title_shape.text = section_title

        return slide

    def save(self, output_path):
        """保存 PPT"""
        self.prs.save(output_path)
        print(f"PPT 已保存到：{output_path}")


def generate_ai_platform_ppt(output_path="ai_middle_platform.pptx"):
    """生成 AI 中台 PPT"""
    generator = AIPPTGenerator(
        title="AI 中台系统",
        subtitle="企业级 AI 能力基础设施"
    )

    # 标题页
    generator.add_title_slide()

    # 目录
    generator.add_content_slide(
        "目录",
        "",
        [
            "1. 什么是 AI 中台",
            "2. 为什么需要 AI 中台",
            "3. 整体架构",
            "4. 核心功能",
            "5. 落地场景",
            "6. 实施路径",
        ]
    )

    # 章节 1
    generator.add_section_slide("什么是 AI 中台")

    generator.add_content_slide(
        "AI 中台定义",
        "企业级 AI 能力的基础设施，将分散的 AI 能力统一纳管、标准化封装、服务化输出",
        [
            "像「水电厂」一样即取即用",
            "像「操作系统」一样屏蔽底层",
            "像「数字神经中枢」一样连接业务与智能",
        ]
    )

    generator.add_content_slide(
        "三个核心理念",
        "",
        [
            "能力复用 — 而非重复建设（一个模型，全行共享）",
            "赋能升级 — 而非推倒重来（2 周完成对接）",
            "持续运营 — 而非项目交付（需要持续运营）",
        ]
    )

    # 章节 2
    generator.add_section_slide("为什么需要 AI 中台")

    generator.add_content_slide(
        "企业痛点",
        "",
        [
            "多系统重复建设 AI 能力",
            "业务系统接入 AI 周期长",
            "数据出域风险、审计缺失",
            "依赖外部厂商、黑盒模型",
            "模型上线后无人运营",
        ]
    )

    # 章节 3
    generator.add_section_slide("整体架构")

    generator.add_content_slide(
        "五层架构",
        "",
        [
            "基础设施层：GPU 集群、信创服务器、容器平台",
            "数据层：向量库、关系库、对象存储、缓存层",
            "AI 中台核心层：模型工厂、知识工厂、智能体工厂",
            "API 网关层：统一入口、认证鉴权、限流熔断",
            "应用层：信贷系统、风险系统、OA 办公、客服系统",
        ]
    )

    # 章节 4
    generator.add_section_slide("核心功能")

    generator.add_content_slide(
        "模型工厂（Model Studio）",
        "",
        [
            "多模型统一接入：一套接口适配所有模型",
            "智能路由：根据请求自动选择最优模型",
            "弹性推理：自动扩缩容应对流量波动",
            "效果监控：实时追踪模型表现",
        ]
    )

    generator.add_content_slide(
        "知识工厂（Knowledge Studio）",
        "",
        [
            "多格式支持：PDF/Word/Excel/PPT/TXT/Markdown",
            "智能分片：按语义而非固定长度切分",
            "混合检索：BM25 + Vector + Rerank",
            "知识管理：权限、版本、增量更新",
        ]
    )

    generator.add_content_slide(
        "智能体工厂（Agent Studio）",
        "",
        [
            "可视化编排：拖拽式工作流设计",
            "记忆管理：短期记忆、长期记忆、共享记忆",
            "工具调用：内置工具、自定义工具",
            "典型场景：智能客服、数据分析、文档处理",
        ]
    )

    # 章节 5
    generator.add_section_slide("落地场景")

    generator.add_content_slide(
        "典型应用场景",
        "",
        [
            "制度文档问答：效率提升 60 倍",
            "合同文本比对：效率提升 5-15 倍",
            "智能客服：解决率 85%+",
            "数据分析报告：分钟级出报告",
            "会议纪要生成：效率提升 10 倍",
        ]
    )

    # 章节 6
    generator.add_section_slide("实施路径")

    generator.add_content_slide(
        "三阶段实施路线",
        "",
        [
            "Phase 1: MVP 版本（3-4 个月）- 验证核心能力",
            "Phase 2: 平台完善（4-6 个月）- 规模化推广",
            "Phase 3: 生态建设（持续运营）- 建立开发者生态",
        ]
    )

    generator.add_content_slide(
        "关键成功因素",
        "",
        [
            "高层支持：AI 中台是一把手工程",
            "场景优先：从高价值场景切入",
            "持续运营：不是项目交付，而是持续运营",
            "生态建设：单靠厂商不够，需要客户/伙伴共建",
        ]
    )

    # 结束页
    generator.add_content_slide(
        "谢谢",
        "如有疑问，欢迎交流",
        []
    )

    generator.save(output_path)
    return output_path


if __name__ == "__main__":
    generate_ai_platform_ppt("AI 中台系统_生成.pptx")
