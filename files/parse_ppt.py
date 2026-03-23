#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT 解析工具 - 解析 PPT 文件结构
"""
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt


def parse_ppt(file_path):
    """解析 PPT 文件"""
    prs = Presentation(file_path)

    slides = []
    for i, slide in enumerate(prs.slides):
        slide_data = {
            "slide_num": i + 1,
            "title": "",
            "content": [],
            "shapes": []
        }

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    if ph_type == 1:  # Title
                        slide_data["title"] = shape.text.strip()
                    else:
                        slide_data["content"].append({
                            "type": "text",
                            "text": shape.text.strip()
                        })
                else:
                    slide_data["shapes"].append({
                        "type": shape.shape_type,
                        "text": shape.text.strip() if hasattr(shape, "text") else ""
                    })

        slides.append(slide_data)

    return {
        "file_name": os.path.basename(file_path),
        "slide_count": len(slides),
        "slides": slides
    }


def save_to_json(data, output_path):
    """保存为 JSON 文件"""
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("用法：python parse_ppt.py <ppt 文件路径> [输出 JSON 路径]")
        sys.exit(1)

    ppt_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else "ppt_structure.json"

    result = parse_ppt(ppt_file)
    save_to_json(result, output_file)
    print(f"PPT 解析完成：{result['slide_count']} 页，已保存到 {output_file}")
