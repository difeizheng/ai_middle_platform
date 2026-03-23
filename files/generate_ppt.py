#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT 生成工具 - 基础版本
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def create_basic_ppt(title, slides_data, output_path):
    """创建基础 PPT"""
    prs = Presentation()

    # 标题页
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = "AI 中台系统"

    # 内容页
    bullet_slide_layout = prs.slide_layouts[1]

    for slide_data in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slide_data.get("title", "")

        tf = body_shape.text_frame
        for i, bullet in enumerate(slide_data.get("bullets", [])):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    prs.save(output_path)
    print(f"PPT 已保存到：{output_path}")


if __name__ == "__main__":
    # 示例用法
    title = "AI 中台系统介绍"
    slides = [
        {"title": "什么是 AI 中台", "bullets": ["企业级 AI 能力基础设施", "一次建设，多次复用"]},
        {"title": "核心价值", "bullets": ["降本增效", "安全合规", "持续运营"]},
        {"title": "架构设计", "bullets": ["五层架构", "七大中枢"]},
    ]

    create_basic_ppt(title, slides, "output.pptx")
