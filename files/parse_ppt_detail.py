#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT 解析工具 - 详细解析
"""
import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def parse_ppt_detail(file_path):
    """详细解析 PPT 文件"""
    prs = Presentation(file_path)

    result = {
        "file_name": os.path.basename(file_path),
        "slide_count": len(prs.slides),
        "slides": []
    }

    for i, slide in enumerate(prs.slides):
        slide_data = {
            "slide_num": i + 1,
            "width": int(prs.slide_width),
            "height": int(prs.slide_height),
            "shapes": []
        }

        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": int(shape.left),
                "top": int(shape.top),
                "width": int(shape.width),
                "height": int(shape.height),
            }

            if hasattr(shape, "text") and shape.text.strip():
                shape_info["text"] = shape.text.strip()

            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                shape_info["table"] = table_data

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_info["image"] = "图片内容"

            slide_data["shapes"].append(shape_info)

        result["slides"].append(slide_data)

    return result


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("用法：python parse_ppt_detail.py <ppt 文件路径>")
        sys.exit(1)

    ppt_file = sys.argv[1]
    result = parse_ppt_detail(ppt_file)
    print(json.dumps(result, ensure_ascii=False, indent=2))
