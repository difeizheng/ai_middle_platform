#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
优化的 PPT 生成工具
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class OptimizedPPTGenerator:
    """优化的 PPT 生成器"""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        # 主题颜色
        self.primary_color = RGBColor(25, 109, 181)  # 蓝色
        self.secondary_color = RGBColor(0, 176, 138)  # 绿色
        self.accent_color = RGBColor(255, 112, 67)  # 橙色

    def add_title_slide(self, title, subtitle):
        """添加标题页"""
        layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)

        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = subtitle

        # 设置标题颜色
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = self.primary_color

        return slide

    def add_section_slide(self, section_num, section_title):
        """添加章节页"""
        layout = self.prs.slide_layouts[6]  # 空白页
        slide = self.prs.slides.add_slide(layout)

        # 添加章节号
        left = Inches(1)
        top = Inches(2)
        width = Inches(10)
        height = Inches(3)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.paragraphs[0]
        p.text = f"PART {section_num}"
        p.font.size = Pt(24)
        p.font.color.rgb = self.secondary_color

        p = tf.add_paragraph()
        p.text = section_title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.primary_color

        return slide

    def add_content_slide(self, title, content_points):
        """添加内容页"""
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        # 设置标题颜色
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.color.rgb = self.primary_color

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame

        for i, point in enumerate(content_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = point
            p.level = 0
            p.font.size = Pt(20)
            p.space_after = Pt(12)

        return slide

    def add_comparison_slide(self, title, columns, data):
        """添加对比表格页"""
        layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        # 创建表格
        rows = len(data) + 1
        cols = len(columns) + 1

        left = Inches(1)
        top = Inches(2)
        width = Inches(11)
        height = Inches(5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # 设置列宽
        for i in range(cols):
            table.columns[i].width = width // cols

        # 填充表头
        for i, col in enumerate(columns):
            table.cell(0, i).text = col
            table.cell(0, i).text_frame.paragraphs[0].font.bold = True

        # 填充数据
        for row_idx, row_data in enumerate(data):
            for col_idx, cell_data in enumerate(row_data):
                table.cell(row_idx + 1, col_idx).text = str(cell_data)

        return slide

    def add_quote_slide(self, quote, author=""):
        """添加金句页"""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)

        left = Inches(2)
        top = Inches(2.5)
        width = Inches(9)
        height = Inches(3)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = quote
        p.font.size = Pt(32)
        p.font.color.rgb = self.primary_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if author:
            p = tf.add_paragraph()
            p.text = f"— {author}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.accent_color
            p.alignment = PP_ALIGN.RIGHT

        return slide

    def save(self, output_path):
        """保存 PPT"""
        self.prs.save(output_path)
        print(f"PPT 已保存到：{output_path}")


def generate_optimized_ppt(output_path="优化版.pptx"):
    """生成优化版 PPT"""
    gen = OptimizedPPTGenerator()

    # 标题页
    gen.add_title_slide(
        "AI 中台解决方案",
        "湖北省农信社"
    )

    # 章节 1
    gen.add_section_slide("01", "什么是 AI 中台")

    gen.add_content_slide(
        "AI 中台定义",
        [
            "企业级 AI 能力的基础设施",
            "将分散的 AI 能力统一纳管、标准化封装、服务化输出",
            "像「水电厂」一样即取即用",
            "像「操作系统」一样屏蔽底层",
            "像「数字神经中枢」一样连接业务与智能",
        ]
    )

    gen.add_content_slide(
        "三个核心理念",
        [
            "能力复用 — 而非重复建设（一个模型，全行共享）",
            "赋能升级 — 而非推倒重来（2 周完成对接）",
            "持续运营 — 而非项目交付（AI 中台是产品）",
        ]
    )

    # 金句页
    gen.add_quote_slide(
        "一次建设，多次复用\n赋能升级，而非推倒重来",
        "AI 中台核心价值"
    )

    # 章节 2
    gen.add_section_slide("02", "为什么需要 AI 中台")

    gen.add_content_slide(
        "企业痛点",
        [
            "多系统重复建设 AI 能力 — 资源浪费",
            "业务系统接入 AI 周期长 — 效率低下",
            "数据出域风险、审计缺失 — 安全隐患",
            "依赖外部厂商、黑盒模型 — 不可控",
            "模型上线后无人运营 — 效果无法保证",
        ]
    )

    # 章节 3
    gen.add_section_slide("03", "整体架构")

    gen.add_content_slide(
        "五层架构",
        [
            "基础设施层：GPU 集群、信创服务器、容器平台",
            "数据层：向量库 (Milvus)、关系库 (PG)、对象存储 (MinIO)",
            "AI 中台核心层：模型工厂、知识工厂、智能体工厂",
            "API 网关层：统一入口、认证鉴权、限流熔断",
            "应用层：信贷、风险、OA、客服等业务系统",
        ]
    )

    gen.add_content_slide(
        "七大中枢",
        [
            "知识中心：统一知识管理",
            "智能体中心：AI 智能体编排",
            "MCP 服务：系统连接器",
            "Skills 技能中心：可复用技能",
            "数据接入：多源数据整合",
            "模型中心：多模型管理",
            "基础保障：安全、监控、运维",
        ]
    )

    # 章节 4
    gen.add_section_slide("04", "核心能力")

    gen.add_content_slide(
        "MCP 协议",
        [
            "Model Context Protocol — 模型上下文协议",
            "标准化系统对接接口",
            "2 周完成现有系统对接",
            "支持数据库、API、消息队列、OA 系统",
        ]
    )

    gen.add_content_slide(
        "全链路审计",
        [
            "AI 黑盒白盒化",
            "每一次交互都可复盘、可定责、可优化",
            "用户登录 → 访问页面 → 调用 API → 查询数据 → 返回结果",
            "完整的审计日志和追踪链路",
        ]
    )

    gen.add_content_slide(
        "私有化部署",
        [
            "数据不出域，满足金融级安全合规",
            "支持信创环境：鲲鹏、海光、昇腾",
            "自主可控，不被单一厂商绑定",
            "灵活扩展，按需部署",
        ]
    )

    # 金句页
    gen.add_quote_slide(
        "赋能现有系统\n而非推倒重来",
        "AI 中台建设理念"
    )

    # 章节 5
    gen.add_section_slide("05", "落地场景")

    gen.add_content_slide(
        "效率提升数据（实测）",
        [
            "项目表单和可研报告审查：3600 倍",
            "投资研究报告编制：2880 倍",
            "供应链效能监测简报：1000 倍+",
            "中标通知书生成：60-240 倍",
            "合同内容比对：5-15 倍",
            "立项前置条件审查：60 倍",
        ]
    )

    gen.add_content_slide(
        "典型应用场景",
        [
            "制度文档问答：智能客服、政策咨询",
            "合同文本比对：风险审查、合规检查",
            "数据分析报告：经营分析、风险监测",
            "会议纪要生成：自动总结、任务提取",
            "智能审批：材料审核、风险评估",
        ]
    )

    # 章节 6
    gen.add_section_slide("06", "实施路径")

    gen.add_content_slide(
        "三阶段实施路线",
        [
            "Phase 1: MVP 版本（3-4 个月）",
            "  - 基础框架、模型工厂、知识工厂",
            "  - 2 个试点场景上线",
            "",
            "Phase 2: 平台完善（4-6 个月）",
            "  - 智能体工厂、MCP 连接器、Skills 市场",
            "  - 5+ 场景规模化推广",
            "",
            "Phase 3: 生态建设（持续运营）",
            "  - 开发者生态、行业解决方案",
        ]
    )

    # 结束页
    gen.add_title_slide(
        "谢谢",
        "如有疑问，欢迎交流"
    )

    gen.save(output_path)
    return output_path


if __name__ == "__main__":
    generate_optimized_ppt("AI 中台解决方案_优化版.pptx")
