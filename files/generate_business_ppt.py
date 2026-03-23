#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
商务推介 PPT 生成工具
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class BusinessPPTGenerator:
    """商务推介 PPT 生成器"""

    def __init__(self, company_name="AI 中台系统"):
        self.company_name = company_name
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, title, subtitle=""):
        """添加标题页"""
        layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        return slide

    def add_content_slide(self, title, bullets):
        """添加内容页"""
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = title

        tf = slide.placeholders[1].text_frame
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

        return slide

    def save(self, output_path):
        """保存 PPT"""
        self.prs.save(output_path)
        print(f"PPT 已保存到：{output_path}")


def generate_business_ppt(output_path="商务推介.pptx"):
    """生成商务推介 PPT"""
    gen = BusinessPPTGenerator()

    # 标题页
    gen.add_title_slide(
        "企业级 AI 中台解决方案",
        "赋能企业智能化转型"
    )

    # 市场机遇
    gen.add_content_slide(
        "AI 市场机遇",
        [
            "2026 年中国 AI 市场规模预计突破 1000 亿元",
            "企业智能化转型需求持续增长",
            "大模型技术成熟，落地应用加速",
            "金融、政府、能源行业需求旺盛",
        ]
    )

    # 产品定位
    gen.add_content_slide(
        "产品定位",
        [
            "企业级 AI 能力基础设施",
            "一次建设，全集团共享",
            "2 周完成现有系统对接",
            "私有化部署，数据不出域",
        ]
    )

    # 核心价值
    gen.add_content_slide(
        "核心价值",
        [
            "降本：避免多系统重复建设",
            "增效：标准化 API，快速接入",
            "合规：全链路审计，满足监管",
            "可控：自主可控，效果可追溯",
        ]
    )

    # 核心功能
    gen.add_content_slide(
        "核心功能",
        [
            "模型工厂：多模型统一接入和管理",
            "知识工厂：文档解析、向量化、检索",
            "智能体工厂：可视化编排 AI 工作流",
            "API 网关：认证鉴权、限流熔断",
        ]
    )

    # 差异化优势
    gen.add_content_slide(
        "差异化优势",
        [
            "MCP 协议：2 周完成系统对接",
            "全链路审计：每一次交互都可复盘",
            "私有化部署：满足金融级安全",
            "持续运营：内置运营体系",
        ]
    )

    # 成功案例
    gen.add_content_slide(
        "成功案例",
        [
            "湖北省农信社：制度问答、合同比对",
            "某大型银行：智能客服、风险预警",
            "某省政府：政务问答、公文处理",
            "某能源企业：设备巡检、安全监控",
        ]
    )

    # 商业模式
    gen.add_content_slide(
        "商业模式",
        [
            "软件授权费（一次性）",
            "实施服务费（一次性）",
            "年度维保费（持续性）",
            "算力服务（可选）",
        ]
    )

    # 实施计划
    gen.add_content_slide(
        "实施计划",
        [
            "Phase 1: MVP 版本（3-4 个月）",
            "Phase 2: 平台完善（4-6 个月）",
            "Phase 3: 生态建设（持续运营）",
        ]
    )

    # 结束页
    gen.add_title_slide(
        "谢谢",
        "期待与您合作"
    )

    gen.save(output_path)
    return output_path


if __name__ == "__main__":
    generate_business_ppt("企业级 AI 中台解决方案_商务推介.pptx")
